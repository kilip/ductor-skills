# /// script
# dependencies = [
#   "python-pptx",
#   "pytesseract",
#   "easyocr",
#   "Pillow",
#   "numpy",
#   "opencv-python-headless",
# ]
# ///

import sys
import os
import io
import re
import tempfile
import subprocess
import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

# Global OCR Engines
TESSERACT_AVAILABLE = False
EASYOCR_READER = None

import pytesseract
import shutil
import platform

# Try to find tesseract in PATH first
tess_path = shutil.which("tesseract")

# Windows specific fallbacks if not in PATH
if not tess_path and platform.system() == "Windows":
    COMMON_TESS_PATHS = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        os.path.expanduser(r"~\AppData\Local\Tesseract-OCR\tesseract.exe"),
    ]
    for path in COMMON_TESS_PATHS:
        if os.path.exists(path):
            tess_path = path
            break

if tess_path:
    pytesseract.pytesseract.tesseract_cmd = tess_path

try:
    pytesseract.get_tesseract_version()
    TESSERACT_AVAILABLE = True
except Exception:
    TESSERACT_AVAILABLE = False

def get_easyocr():
    global EASYOCR_READER
    if EASYOCR_READER is None:
        import easyocr
        EASYOCR_READER = easyocr.Reader(['en', 'id'], gpu=False)
    return EASYOCR_READER

def download_from_gdrive(file_id_or_url, account=None):
    """
    Downloads a file from Google Drive using the 'gog' CLI tool.
    """
    # Extract file ID if a URL is provided
    file_id = file_id_or_url
    if "drive.google.com" in file_id_or_url:
        match = re.search(r"/d/([-\w]{25,})", file_id_or_url)
        if match:
            file_id = match.group(1)
    
    temp_dir = os.path.join(tempfile.gettempdir(), "slide-extractor")
    os.makedirs(temp_dir, exist_ok=True)
    
    dest_path = os.path.join(temp_dir, f"{file_id}.pptx")
    
    # If file already exists in temp, we can reuse it or redownload.
    # Let's redownload to be sure it's fresh, unless Pak Bos prefers otherwise.
    # But for efficiency, let's check if it exists first.
    if os.path.exists(dest_path):
        print(f"INFO: Using cached file at {dest_path}")
        return dest_path

    print(f"INFO: Downloading file {file_id} from Google Drive...")
    try:
        # Build command: gog drive download <file_id> --out <dest_path> [--account <account>]
        cmd = ["gog", "drive", "download", file_id, "--out", dest_path]
        if account:
            cmd.extend(["--account", account])
            
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True
        )
        if result.returncode != 0:
            print(f"Error downloading from GDrive: {result.stderr}")
            return None
        
        if not os.path.exists(dest_path):
            print("Error: gog finished but file was not found at expected path.")
            return None
            
        return dest_path
    except FileNotFoundError:
        print("Error: 'gog' CLI tool not found. Please install it to use GDrive features.")
        return None
    except Exception as e:
        print(f"Error: {e}")
        return None

def extract_text_from_pptx(pptx_path, preferred_engine=None):
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}")
        return

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error opening PPTX: {e}")
        return

    def perform_ocr(img_blob):
        # Priority 1: Tesseract
        if (preferred_engine is None or preferred_engine == "tesseract") and TESSERACT_AVAILABLE:
            try:
                img = Image.open(io.BytesIO(img_blob))
                ocr_text = pytesseract.image_to_string(img).strip()
                if ocr_text:
                    return ocr_text
            except Exception:
                pass

        # Priority 2: EasyOCR
        if preferred_engine is None or preferred_engine == "easyocr":
            try:
                reader = get_easyocr()
                img = Image.open(io.BytesIO(img_blob)).convert("RGB")
                img_array = np.array(img)
                results = reader.readtext(img_array, detail=0)
                ocr_text = " ".join(results).strip()
                if ocr_text:
                    return ocr_text
            except Exception:
                pass
        
        return None

    def get_text_and_images(shape, text_runs):
        if hasattr(shape, "text") and shape.text.strip():
            text_runs.append(shape.text.strip())
        
        if shape.has_table:
            for row in shape.table.rows:
                row_text = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()]
                if row_text:
                    text_runs.append(" | ".join(row_text))

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                get_text_and_images(s, text_runs)

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ocr_res = perform_ocr(shape.image.blob)
            if ocr_res:
                text_runs.append(ocr_res)

    # Metadata Header
    print("---")
    print(f"title: {os.path.splitext(os.path.basename(pptx_path))[0]}")
    print(f"filename: {os.path.basename(pptx_path)}")
    print(f"slide_count: {len(prs.slides)}")
    print("---")
    print()

    for i, slide in enumerate(prs.slides):
        print(f"# Slide {i + 1}")
        text_runs = []
        for shape in slide.shapes:
            get_text_and_images(shape, text_runs)

        if text_runs:
            unique_text = []
            seen = set()
            for t in text_runs:
                clean_t = t.strip().lower()
                if clean_t and clean_t not in seen:
                    unique_text.append(t)
                    seen.add(clean_t)
            print("\n".join(unique_text))
        else:
            print("*[Empty Slide]*")
        print("\n---\n")

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Extract text from PPTX with OCR fallback.")
    parser.add_argument("pptx_path", help="Path to the PPTX file or Google Drive File ID/URL")
    parser.add_argument("--engine", choices=["tesseract", "easyocr"], help="Force specific OCR engine")
    parser.add_argument("--account", "-a", help="Google Drive account name for 'gog' CLI")
    args = parser.parse_args()

    target_path = args.pptx_path
    
    # If path doesn't exist locally, check if it's a GDrive ID/URL
    if not os.path.exists(target_path):
        # Basic check: if it looks like an ID or contains drive.google.com
        if "drive.google.com" in target_path or re.match(r"^[-\w]{25,}$", target_path):
            downloaded_path = download_from_gdrive(target_path, account=args.account)
            if downloaded_path:
                target_path = downloaded_path
            else:
                sys.exit(1)
        else:
            print(f"Error: File not found at {target_path}")
            sys.exit(1)

    extract_text_from_pptx(target_path, args.engine)
