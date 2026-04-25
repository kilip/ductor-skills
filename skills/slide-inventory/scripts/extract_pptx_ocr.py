# /// script
# dependencies = [
#   "python-pptx",
#   "pytesseract",
#   "Pillow",
# ]
# ///

import os
import sys
import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pytesseract
from PIL import Image
import io

def extract_text_and_ocr(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"Error: File {pptx_path} not found.")
        return

    prs = Presentation(pptx_path)
    full_output = []

    for i, slide in enumerate(prs.slides):
        slide_text = [f"--- Slide {i+1} ---"]
        
        # 1. Direct text extraction
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            
            # 2. OCR on images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_bytes = shape.image.blob
                image = Image.open(io.BytesIO(image_bytes))
                try:
                    ocr_text = pytesseract.image_to_string(image, lang='ind+eng')
                    if ocr_text.strip():
                        slide_text.append("[OCR Content]:")
                        slide_text.append(ocr_text.strip())
                except Exception as e:
                    slide_text.append(f"[OCR Error]: {str(e)}")
            
            # 3. OCR on grouped shapes or other types that might contain images
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for subshape in shape.shapes:
                    if subshape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_bytes = subshape.image.blob
                        image = Image.open(io.BytesIO(image_bytes))
                        try:
                            ocr_text = pytesseract.image_to_string(image, lang='ind+eng')
                            if ocr_text.strip():
                                slide_text.append("[OCR Content (Grouped)]: ")
                                slide_text.append(ocr_text.strip())
                        except Exception as e:
                            pass

        full_output.append("\n".join(slide_text))

    return "\n\n".join(full_output)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract text and perform OCR on PPTX files.")
    parser.add_argument("pptx", help="Path to the PPTX file")
    parser.add_argument("--output", help="Path to save the output text file", default=None)
    
    args = parser.parse_args()
    
    result = extract_text_and_ocr(args.pptx)
    
    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(result)
        print(f"Output saved to {args.output}")
    else:
        print(result)
